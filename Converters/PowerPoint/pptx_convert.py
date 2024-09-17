from pptx import Presentation
from Utils.pdf_utils import setup_pdf, add_text_to_pdf

def convert_pptx_to_pdf(pptx_path, pdf_output_path):

    prs = Presentation(pptx_path)
    c = setup_pdf(pdf_output_path)

    y_offset = 750
    for slide in prs.slides:
        slide_text = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        y_offset = add_text_to_pdf(c, slide_text, y_offset=y_offset)
        c.showPage()  # New page per slide

    c.save()
    print(f"Successfully converted {pptx_path} to {pdf_output_path}")
